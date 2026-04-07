"""ResearchPublisher — writes research artifacts to docs/ and commits to GitHub.

Pipeline (called as the "publish" stage of the research workflow):
  1. Parse `### File: <name>` blocks from research_report + content_creator output
  2. Write each file to docs/research/REQ-XXX-<slug>/
  3. Render slides.md → slides.pptx (python-pptx)
  4. Render report.md → report.pdf (weasyprint)
  5. Commit all files to GitHub via the Trees API in one atomic commit
  6. Soft-fail: if GitHub publish errors, the request still completes; the local
     files remain on disk and the error is captured in the result.
"""

import re
from pathlib import Path
from typing import Any

import structlog
from slugify import slugify

from src.core.github_publisher import GitHubPublishError, GitHubPublisher

logger = structlog.get_logger()

# Where artifacts are written, relative to the project root inside the container
DOCS_RESEARCH_DIR = Path("docs/research")


class ResearchPublisher:
    """Parses research artifacts, writes them locally, and publishes to GitHub."""

    def __init__(self, project_root: str = ".") -> None:
        self.root = Path(project_root)
        self.docs_dir = self.root / DOCS_RESEARCH_DIR
        self.github = GitHubPublisher()

    async def publish(
        self, request_id: str, description: str, artifacts: dict[str, Any]
    ) -> dict[str, Any]:
        """Main entry point — called by the workflow runner for the publish stage.

        Args:
            request_id: The REQ-XXX id
            description: The original request description (used for the folder slug)
            artifacts: workflow artifacts dict — contains research_specialist_output
                       and content_creator_output (or similar keys)

        Returns:
            dict with keys: published_files, commit_sha, commit_url, publish_error
        """
        result: dict[str, Any] = {
            "published_files": [],
            "commit_sha": None,
            "commit_url": None,
            "publish_error": None,
        }

        # ── 1. Collect text from both agents ──────
        research_text = self._extract_agent_output(artifacts, "research_specialist")
        content_text = self._extract_agent_output(artifacts, "content_creator")

        if not research_text and not content_text:
            result["publish_error"] = "No research_specialist or content_creator output found"
            logger.warning("research_publish_no_input", request_id=request_id)
            return result

        # ── 2. Determine the folder name ──────────
        folder_name = self._make_folder_name(request_id, description)
        folder_path = self.docs_dir / folder_name
        folder_path.mkdir(parents=True, exist_ok=True)
        logger.info("research_publish_folder", path=str(folder_path), request_id=request_id)

        # ── 3. Parse file blocks from content_creator output ──────────
        files_written: dict[str, str] = {}  # rel_path → text content

        # Always include the raw research report as research-report.md
        if research_text:
            files_written["research-report.md"] = research_text

        if content_text:
            parsed = self._parse_file_blocks(content_text)
            files_written.update(parsed)

        if not files_written:
            result["publish_error"] = "No file blocks found in content_creator output"
            logger.warning("research_publish_no_files", request_id=request_id)
            return result

        # ── 4. Write text files to disk ──────────
        for filename, content in files_written.items():
            file_path = folder_path / filename
            file_path.parent.mkdir(parents=True, exist_ok=True)
            file_path.write_text(content, encoding="utf-8")
            logger.debug("research_file_written", file=filename, size=len(content))

        # ── 5. Render derived artifacts (PDF, PPTX) ──────────
        binary_files: dict[str, bytes] = {}  # rel_path → bytes content

        if "report.md" in files_written:
            try:
                pdf_bytes = self._render_pdf(files_written["report.md"])
                (folder_path / "report.pdf").write_bytes(pdf_bytes)
                binary_files["report.pdf"] = pdf_bytes
                logger.info("research_pdf_rendered", request_id=request_id, size=len(pdf_bytes))
            except Exception as e:
                logger.warning("research_pdf_render_failed", error=str(e), request_id=request_id)

        if "slides.md" in files_written:
            try:
                pptx_bytes = self._render_pptx(files_written["slides.md"])
                (folder_path / "slides.pptx").write_bytes(pptx_bytes)
                binary_files["slides.pptx"] = pptx_bytes
                logger.info("research_pptx_rendered", request_id=request_id, size=len(pptx_bytes))
            except Exception as e:
                logger.warning("research_pptx_render_failed", error=str(e), request_id=request_id)

        # Collect all written files (text + binary)
        all_files = list(files_written.keys()) + list(binary_files.keys())
        result["published_files"] = [f"{folder_name}/{f}" for f in all_files]

        # ── 6. Publish to GitHub via Trees API (soft-fail) ──────────
        if not self.github.is_configured():
            result["publish_error"] = "GITHUB_TOKEN/GITHUB_REPO not configured — files written locally only"
            logger.warning("research_publish_not_configured", request_id=request_id)
            return result

        # Build the {repo_path: content} dict for the Trees API
        commit_files: dict[str, bytes | str] = {}
        for filename, content in files_written.items():
            commit_files[f"docs/research/{folder_name}/{filename}"] = content
        for filename, content_bytes in binary_files.items():
            commit_files[f"docs/research/{folder_name}/{filename}"] = content_bytes

        commit_msg = (
            f"docs({request_id}): research artifacts for {description[:60]}\n\n"
            f"Auto-published by Agent Team Research Pipeline.\n"
            f"Folder: docs/research/{folder_name}/\n"
            f"Files: {len(commit_files)}"
        )

        try:
            commit_info = await self.github.commit_files(commit_files, commit_msg)
            result["commit_sha"] = commit_info["short_sha"]
            result["commit_url"] = commit_info["url"]
            logger.info(
                "research_published_to_github",
                request_id=request_id,
                sha=commit_info["short_sha"],
                files=len(all_files),
            )
        except GitHubPublishError as e:
            result["publish_error"] = str(e)
            logger.warning("research_github_publish_failed", request_id=request_id, error=str(e))

        return result

    # ── Helpers ──────────────────────────────────────────────────────

    def _extract_agent_output(self, artifacts: dict[str, Any], agent_id: str) -> str:
        """Pull an agent's output text from the workflow artifacts dict.

        The runner stores outputs under keys like `<agent_id>_output` or as the
        agent's named output (e.g., `research_report`, `research_artifacts`).
        """
        for key in (f"{agent_id}_output", agent_id):
            val = artifacts.get(key)
            if isinstance(val, str) and val.strip():
                return val
        # Fallback: scan for any string artifact that mentions the agent
        for key, val in artifacts.items():
            if agent_id in key and isinstance(val, str) and val.strip():
                return val
        return ""

    def _make_folder_name(self, request_id: str, description: str) -> str:
        """REQ-A3F2C1 + 'Compare vector databases' → 'REQ-A3F2C1-compare-vector-databases'."""
        # Strip markdown attachments section if any
        clean = description.split("\n\n**Attachments:**")[0].strip()
        # Take the first line / sentence and slug it, max 60 chars
        first_chunk = re.split(r"[.\n]", clean)[0][:80]
        slug = slugify(first_chunk, max_length=60, word_boundary=True)
        if not slug:
            slug = "research"
        return f"{request_id}-{slug}"

    def _parse_file_blocks(self, text: str) -> dict[str, str]:
        r"""Parse ``### File: `filename` `` blocks followed by fenced code blocks."""
        files: dict[str, str] = {}
        # Match: ### File: `filename`\n```lang\n<content>\n```
        pattern = r"###\s+File:\s*`([^`]+)`\s*\n```[\w]*\n([\s\S]*?)\n```"
        for match in re.finditer(pattern, text):
            filename = match.group(1).strip()
            content = match.group(2).strip()
            # Security: no path traversal
            if ".." in filename or filename.startswith("/"):
                logger.warning("research_path_traversal_blocked", filename=filename)
                continue
            if filename and content:
                files[filename] = content + "\n"
        return files

    def _render_pdf(self, markdown_text: str) -> bytes:
        """Convert markdown → HTML → PDF using weasyprint."""
        import markdown as md
        from weasyprint import HTML

        html_body = md.markdown(markdown_text, extensions=["tables", "fenced_code", "toc"])
        html_doc = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<style>
  body {{ font-family: 'DejaVu Sans', sans-serif; max-width: 720px; margin: 40px auto; line-height: 1.5; color: #1f2937; }}
  h1 {{ color: #111827; border-bottom: 2px solid #2563eb; padding-bottom: 6px; }}
  h2 {{ color: #1f2937; margin-top: 24px; }}
  h3 {{ color: #374151; }}
  table {{ border-collapse: collapse; width: 100%; margin: 12px 0; }}
  th, td {{ border: 1px solid #e5e7eb; padding: 6px 10px; text-align: left; }}
  th {{ background: #f3f4f6; }}
  code {{ background: #f3f4f6; padding: 1px 4px; border-radius: 3px; font-family: 'DejaVu Sans Mono', monospace; }}
  pre {{ background: #f3f4f6; padding: 12px; border-radius: 6px; overflow-x: auto; }}
  blockquote {{ border-left: 3px solid #2563eb; padding-left: 12px; color: #4b5563; }}
</style>
</head>
<body>
{html_body}
</body>
</html>"""
        return HTML(string=html_doc).write_pdf()

    def _render_pptx(self, slides_markdown: str) -> bytes:
        """Convert slides markdown → PowerPoint deck using python-pptx.

        Format expected: each slide separated by `---`. First line is the title
        (with or without `#`). Bullets start with `-` or `*`. Lines starting with
        `Speaker notes:` go into speaker notes.
        """
        import io
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Split on lines that are exactly `---`
        slide_blocks = re.split(r"\n\s*---+\s*\n", slides_markdown)

        for block in slide_blocks:
            block = block.strip()
            if not block:
                continue

            lines = [l.rstrip() for l in block.split("\n") if l.strip()]
            if not lines:
                continue

            # Extract title (first non-empty line, strip leading #'s)
            title = re.sub(r"^#+\s*", "", lines[0]).strip()
            # Strip "Slide N:" prefix if present
            title = re.sub(r"^Slide\s+\d+:\s*", "", title, flags=re.IGNORECASE).strip()

            # Bullets and speaker notes
            bullets: list[str] = []
            speaker_notes_parts: list[str] = []
            in_notes = False
            for line in lines[1:]:
                if line.lower().startswith("speaker notes:"):
                    in_notes = True
                    note_text = line.split(":", 1)[1].strip()
                    if note_text:
                        speaker_notes_parts.append(note_text)
                    continue
                if in_notes:
                    speaker_notes_parts.append(line)
                else:
                    cleaned = re.sub(r"^[-*]\s*", "", line).strip()
                    if cleaned:
                        bullets.append(cleaned)

            # Use a Title+Content layout
            layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title or "Untitled"

            # Add bullets to the content placeholder
            if bullets and len(slide.placeholders) > 1:
                content_ph = slide.placeholders[1]
                tf = content_ph.text_frame
                tf.text = bullets[0]
                for bullet in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(18)

            # Add speaker notes
            if speaker_notes_parts:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = " ".join(speaker_notes_parts)

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

